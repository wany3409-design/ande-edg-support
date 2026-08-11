"""
PPTX 解析器

基于 python-pptx 提取文本，支持：
- 提取幻灯片文本
- 保留标题信息
- 提取表格内容
"""

from typing import List

from pptx import Presentation

from src.document_parser.base_parser import BaseDocumentParser
from src.document_parser.models import RawDocument, ParsedPage


class PptxParser(BaseDocumentParser):
    """PPTX 演示文稿解析器"""

    supported_extensions = [".pptx"]
    parser_name = "pptx"

    def _do_parse(self, raw_doc: RawDocument) -> List[ParsedPage]:
        try:
            prs = Presentation(raw_doc.file_path)
        except Exception as e:
            raise RuntimeError(f"无法打开PPTX文件: {raw_doc.file_path}, 错误: {e}")

        pages = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            section_hint = ""

            for shape in slide.shapes:
                # 提取文本框
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)

                # 提取表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)

            if texts:
                # 第一个非空文本作为章节提示
                if texts and len(texts[0]) < 100:
                    section_hint = texts[0]

                pages.append(ParsedPage(
                    page_num=slide_num,
                    text="\n".join(texts),
                    section=section_hint,
                ))

        return pages
