"""
技术支持助手 — Streamlit UI

直接调用 AnswerPipeline，不经过 FastAPI。
启动: streamlit run src/ui/app.py
"""

import sys
import io
import time
import json
import logging
from datetime import datetime, timezone, timedelta
from pathlib import Path

_PROJECT_ROOT = Path(__file__).parent.parent.parent
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

import streamlit as st

# ===== 日志 =====
_LOG_DIR = _PROJECT_ROOT / "data"
_LOG_DIR.mkdir(parents=True, exist_ok=True)
_LOGGER = logging.getLogger("tech_support_ui")
_LOGGER.setLevel(logging.INFO)
_h = logging.FileHandler(_LOG_DIR / "app.log", encoding="utf-8")
_h.setFormatter(logging.Formatter("%(asctime)s | %(levelname)s | %(message)s"))
_LOGGER.addHandler(_h)
_LOGGER.info("UI session started")
_TZ_BEIJING = timezone(timedelta(hours=8))


def _now():
    return datetime.now(_TZ_BEIJING).isoformat(timespec="seconds")


# ===== 页面配置 =====
st.set_page_config(
    page_title="技术支持助手",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ===== 全局 CSS =====
st.markdown("""
<style>
    /* ========== 导入字体 ========== */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=Noto+Sans+SC:wght@400;500;700&display=swap');

    :root {
        --primary: #2563eb;
        --primary-light: #dbeafe;
        --primary-dark: #1e40af;
        --gray-50: #f9fafb;
        --gray-100: #f3f4f6;
        --gray-200: #e5e7eb;
        --gray-300: #d1d5db;
        --gray-500: #6b7280;
        --gray-600: #4b5563;
        --gray-700: #374151;
        --gray-800: #1f2937;
        --gray-900: #111827;
    }

    html, body, [class*="css"] {
        font-family: 'Inter', 'Noto Sans SC', -apple-system, BlinkMacSystemFont, sans-serif;
    }

    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }
    header[data-testid="stHeader"] { background: transparent !important; }

    .main .block-container {
        padding-top: 1.5rem;
        max-width: 900px;
    }

    /* ========== 顶部 Header ========== */
    .app-header {
        display: flex;
        align-items: center;
        justify-content: center;
        gap: 14px;
        padding: 16px 20px;
        background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%);
        border-radius: 14px;
        margin-bottom: 6px;
        box-shadow: 0 4px 12px rgba(37,99,235,0.18);
        text-align: center;
    }
    .app-header-icon {
        width: 44px; height: 44px;
        background: rgba(255,255,255,0.18);
        border-radius: 10px;
        display: flex; align-items: center; justify-content: center;
        font-size: 24px;
        flex-shrink: 0;
    }
    .app-header-text h1 {
        color: #fff;
        font-size: 1.35rem;
        font-weight: 700;
        margin: 0;
        letter-spacing: -0.01em;
        text-align: center;
    }
    .app-header-text p {
        color: rgba(255,255,255,0.78);
        font-size: 0.82rem;
        margin: 2px 0 0 0;
        font-weight: 400;
        text-align: center;
    }

    /* ========== 用户消息气泡 ========== */
    .user-msg-wrapper {
        display: flex;
        justify-content: flex-end;
        margin: 12px 0;
    }
    .user-bubble {
        background: linear-gradient(135deg, #2563eb, #3b82f6);
        color: #fff;
        padding: 10px 16px;
        border-radius: 16px 16px 4px 16px;
        max-width: 75%;
        font-size: 0.93rem;
        line-height: 1.55;
        box-shadow: 0 2px 8px rgba(37,99,235,0.25);
        word-break: break-word;
    }
    .user-bubble .user-label {
        font-size: 0.7rem;
        opacity: 0.75;
        margin-bottom: 4px;
        letter-spacing: 0.02em;
    }
    .user-attach {
        display: inline-flex; align-items: center; gap: 4px;
        margin-top: 8px;
        padding: 3px 10px;
        background: rgba(255,255,255,0.18);
        border-radius: 12px;
        font-size: 0.76rem;
    }

    /* ========== AI 回答卡片 ========== */
    .ai-card {
        background: #fff;
        border: 1px solid var(--gray-200);
        border-radius: 14px;
        box-shadow: 0 1px 3px rgba(0,0,0,0.08), 0 1px 2px rgba(0,0,0,0.06);
        margin: 12px 0 18px 0;
        overflow: hidden;
    }
    .ai-card-header {
        display: flex;
        align-items: center;
        gap: 8px;
        padding: 12px 18px;
        border-bottom: 1px solid var(--gray-100);
        background: var(--gray-50);
    }
    .ai-card-header .badge {
        display: inline-flex;
        align-items: center;
        gap: 4px;
        padding: 4px 10px;
        border-radius: 20px;
        font-size: 0.78rem;
        font-weight: 600;
    }
    .badge-knowledge { background: #dbeafe; color: #1e40af; }
    .badge-implementation { background: #d1fae5; color: #065f46; }
    .badge-troubleshooting { background: #fef3c7; color: #92400e; }
    .badge-refusal { background: #fee2e2; color: #991b1b; }
    .ai-card-body {
        padding: 18px 20px;
        font-size: 0.92rem;
        line-height: 1.7;
        color: var(--gray-800);
    }
    .ai-card-body h3 {
        font-size: 1.05rem;
        font-weight: 700;
        color: var(--gray-900);
        margin: 16px 0 8px 0;
        padding-bottom: 6px;
        border-bottom: 2px solid var(--gray-100);
    }
    .ai-card-body h3:first-child { margin-top: 0; }
    .ai-card-body strong { color: var(--gray-900); }
    .ai-card-body ul, .ai-card-body ol { padding-left: 20px; margin: 6px 0; }
    .ai-card-body li { margin: 4px 0; }
    .ai-card-body blockquote {
        border-left: 3px solid var(--primary);
        background: var(--primary-light);
        padding: 10px 14px;
        border-radius: 0 6px 6px 0;
        margin: 10px 0;
        color: var(--gray-700);
    }

    /* ========== 证据分级标记高亮 ========== */
    .tag-supported { background: #d1fae5; color: #065f46; padding: 1px 7px; border-radius: 5px; font-size: 0.82em; font-weight: 600; white-space: nowrap; }
    .tag-integrated { background: #e0e7ff; color: #3730a3; padding: 1px 7px; border-radius: 5px; font-size: 0.82em; font-weight: 600; white-space: nowrap; }
    .tag-inferred { background: #fef3c7; color: #92400e; padding: 1px 7px; border-radius: 5px; font-size: 0.82em; font-weight: 600; white-space: nowrap; }
    .tag-missing { background: #fee2e2; color: #991b1b; padding: 1px 7px; border-radius: 5px; font-size: 0.82em; font-weight: 600; white-space: nowrap; }
    .tag-generic { background: #f3e8ff; color: #6b21a8; padding: 1px 7px; border-radius: 5px; font-size: 0.82em; font-weight: 600; white-space: nowrap; }

    /* ========== 拒绝回答卡片 ========== */
    .ai-card-refusal { border-color: #fecaca; }
    .ai-card-refusal .ai-card-header { background: #fef2f2; border-bottom-color: #fee2e2; }
    .ai-card-refusal .ai-card-body { color: var(--gray-600); }

    /* ========== 参考资料 ========== */
    .ref-section {
        margin: 14px 0;
        border: 1px solid var(--gray-200);
        border-radius: 10px;
        overflow: hidden;
        box-shadow: 0 1px 2px rgba(0,0,0,0.05);
    }
    .ref-item {
        display: flex;
        gap: 12px;
        padding: 12px 16px;
        border-bottom: 1px solid var(--gray-100);
        align-items: flex-start;
        transition: background 0.12s;
    }
    .ref-item:hover { background: var(--gray-50); }
    .ref-item:last-child { border-bottom: none; }
    .ref-index {
        width: 26px; height: 26px;
        background: var(--primary-light);
        color: var(--primary-dark);
        border-radius: 50%;
        display: flex; align-items: center; justify-content: center;
        font-size: 0.75rem;
        font-weight: 700;
        flex-shrink: 0;
        margin-top: 2px;
    }
    .ref-info { flex: 1; min-width: 0; }
    .ref-info .ref-src { font-weight: 600; font-size: 0.88rem; color: var(--gray-800); }
    .ref-info .ref-meta { font-size: 0.78rem; color: var(--gray-500); margin-top: 2px; }
    .ref-info .ref-meta span { display: inline-block; margin-right: 12px; }
    .ref-evidence-text {
        background: #f8fafc;
        color: var(--gray-700);
        font-size: 0.82rem;
        line-height: 1.6;
        padding: 10px 14px;
        border-radius: 6px;
        margin-top: 8px;
        white-space: pre-wrap;
        word-break: break-word;
        border-left: 3px solid var(--gray-300);
    }

    /* ========== 检索详情 ========== */
    .retrieval-item {
        padding: 12px 14px;
        border-bottom: 1px solid var(--gray-100);
        transition: background 0.12s;
    }
    .retrieval-item:hover { background: var(--gray-50); }
    .retrieval-item:last-child { border-bottom: none; }
    .score-tag {
        display: inline-block;
        padding: 3px 8px;
        border-radius: 4px;
        font-weight: 600;
        font-size: 0.75rem;
    }
    .score-high { background: #d1fae5; color: #065f46; }
    .score-mid { background: #fef3c7; color: #92400e; }
    .score-low { background: #fee2e2; color: #991b1b; }

    /* ========== Timeline 指标 ========== */
    .metrics-row { display: flex; gap: 10px; margin: 10px 0; flex-wrap: wrap; }
    .metric-chip {
        display: inline-flex; align-items: center; gap: 5px;
        padding: 6px 12px;
        background: var(--gray-50);
        border-radius: 20px;
        font-size: 0.78rem;
        color: var(--gray-600);
        border: 1px solid var(--gray-200);
    }
    .metric-chip .val { font-weight: 600; color: var(--gray-800); }

    /* ========== 按钮 ========== */
    div[data-testid="stButton"] > button {
        border-radius: 20px !important;
        font-size: 0.82rem !important;
        font-weight: 500 !important;
        padding: 6px 16px !important;
        transition: all 0.15s !important;
    }
    div[data-testid="stButton"] > button:hover {
        transform: translateY(-1px);
        box-shadow: 0 4px 6px rgba(0,0,0,0.06);
    }

    /* ========== Grounding 警告 ========== */
    .grounding-warn {
        display: flex;
        align-items: flex-start;
        gap: 8px;
        padding: 10px 14px;
        background: #fffbeb;
        border: 1px solid #fde68a;
        border-radius: 6px;
        font-size: 0.82rem;
        color: #92400e;
        margin: 10px 0;
    }

    /* ========== 聊天输入框 ========== */
    [data-testid="stChatInput"] textarea {
        border-radius: 24px !important;
        border: 2px solid var(--gray-200) !important;
        padding: 12px 20px !important;
        font-size: 0.93rem !important;
        transition: border-color 0.2s !important;
    }
    [data-testid="stChatInput"] textarea:focus {
        border-color: var(--primary) !important;
        box-shadow: 0 0 0 3px rgba(37,99,235,0.1) !important;
    }

    /* ========== 空状态 ========== */
    .empty-state {
        text-align: center;
        padding: 60px 20px;
        color: var(--gray-400);
    }
    .empty-state .icon { font-size: 48px; margin-bottom: 16px; }
    .empty-state .title { font-size: 1.1rem; font-weight: 600; color: var(--gray-500); margin-bottom: 8px; }
    .empty-state .hints { font-size: 0.85rem; line-height: 1.8; }
    .empty-state .hint-chip {
        display: inline-block;
        padding: 4px 12px;
        background: var(--gray-100);
        border-radius: 16px;
        margin: 3px;
        cursor: default;
        color: var(--gray-600);
        font-size: 0.8rem;
    }

    /* ========== 响应式 ========== */
    @media (max-width: 640px) {
        .user-bubble { max-width: 90%; }
        .app-header { flex-direction: column; text-align: center; }
        .app-header-text h1 { font-size: 1.15rem; }
    }

    /* ========== 附件按钮紧贴底部输入框 ========== */
    .st-key-attach_toggle {
        margin-bottom: -0.5rem !important;
    }
    [data-testid="stChatInput"] {
        margin-top: -1rem !important;
    }
</style>
""", unsafe_allow_html=True)


# ===== 证据分级标记 → 高亮样式 =====
_EVIDENCE_TAG_STYLES = {
    "【资料明确说明】": "tag-supported",
    "【跨文档整理】": "tag-integrated",
    "【根据多份资料整理】": "tag-integrated",
    "【推断】": "tag-inferred",
    "【资料未说明】": "tag-missing",
    "【资料无法确认】": "tag-missing",
    "【非安得产品资料】": "tag-generic",
}


# ===== Session State =====
def _init_session():
    defaults = {
        "chat_history": [],
        "chat_messages": [],
        "last_evidences": None,
        "last_query": None,
        "last_result": None,
        "feedback_given": {},
        "msg_counter": 0,
        "upload_text": "",
        "upload_meta": None,
        "show_attach": False,
    }
    for key, val in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = val


_init_session()


# ===== 管线 =====
@st.cache_resource(show_spinner=False)
def _load_pipeline():
    from src.llm.answer_pipeline import AnswerPipeline
    _LOGGER.info("Initializing AnswerPipeline...")
    try:
        p = AnswerPipeline()
        _ = p.embedding_model
        _ = p.collection
        _LOGGER.info(f"Pipeline ready, collection={p.collection.count()}")
        return p
    except Exception as e:
        _LOGGER.error(f"Pipeline init failed: {e}", exc_info=True)
        raise


# ===== 上传文件文本提取 =====
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".ico"}
_TEXT_EXTS = {".txt", ".md", ".log", ".csv", ".json", ".yaml", ".yml", ".xml", ".py", ".sh", ".ini", ".conf"}
_UPLOAD_TYPES = ["pdf", "docx", "pptx", "xlsx", "xls", "txt", "md", "log", "csv",
                 "json", "html", "htm", "png", "jpg", "jpeg", "webp"]


def _extract_text(uploaded_file) -> tuple:
    """从上传文件提取文本。返回 (text, kind)，kind ∈ {"text", "image", "unsupported"}"""
    name = (uploaded_file.name or "").lower()
    ext = Path(name).suffix
    data = uploaded_file.getvalue()

    if ext in _IMAGE_EXTS:
        return "", "image"

    if ext in _TEXT_EXTS:
        for enc in ("utf-8", "gbk"):
            try:
                return data.decode(enc), "text"
            except (UnicodeDecodeError, LookupError):
                continue
        return data.decode("utf-8", errors="ignore"), "text"

    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(stream=data, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text, "text" if text.strip() else "unsupported"
        except Exception as e:
            _LOGGER.warning(f"PDF extract failed: {e}")
            return "", "unsupported"

    if ext == ".docx":
        try:
            import docx
            d = docx.Document(io.BytesIO(data))
            text = "\n".join(p.text for p in d.paragraphs)
            return text, "text" if text.strip() else "unsupported"
        except Exception as e:
            _LOGGER.warning(f"DOCX extract failed: {e}")
            return "", "unsupported"

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text)
            text = "\n".join(parts)
            return text, "text" if text.strip() else "unsupported"
        except Exception as e:
            _LOGGER.warning(f"PPTX extract failed: {e}")
            return "", "unsupported"

    if ext in {".xlsx", ".xlsm"}:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
            wb.close()
            text = "\n".join(parts)
            return text, "text" if text.strip() else "unsupported"
        except Exception as e:
            _LOGGER.warning(f"XLSX extract failed: {e}")
            return "", "unsupported"

    if ext in {".html", ".htm"}:
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(data, "html.parser")
            text = soup.get_text("\n")
            return text, "text" if text.strip() else "unsupported"
        except Exception as e:
            _LOGGER.warning(f"HTML extract failed: {e}")
            return "", "unsupported"

    for enc in ("utf-8", "gbk"):
        try:
            return data.decode(enc), "text"
        except (UnicodeDecodeError, LookupError):
            continue
    return "", "unsupported"


# ===== 日志工具 =====
def _log_query(query, result, wall_time):
    try:
        entry = {
            "timestamp": _now(), "query": query,
            "category": result.get("category", ""),
            "is_answerable": result.get("is_answerable", False),
            "retrieval_count": result.get("retrieval_count", 0),
            "rerank_count": result.get("rerank_count", 0),
            "top1_sim": (result.get("evidences", [{}])[0] or {}).get("similarity", 0) if result.get("evidences") else 0,
            "top1_rerank": (result.get("evidences", [{}])[0] or {}).get("rerank_score", 0) if result.get("evidences") else 0,
            "wall_time": round(wall_time, 1),
        }
        _LOGGER.info(f"QUERY | {json.dumps(entry, ensure_ascii=False)}")
    except Exception:
        pass


def _log_feedback(query, answer, feedback):
    try:
        entry = {"timestamp": _now(), "query": query[:500], "answer": answer[:500], "feedback": feedback}
        with open(_LOG_DIR / "feedback.jsonl", "a", encoding="utf-8") as f:
            f.write(json.dumps(entry, ensure_ascii=False) + "\n")
        _LOGGER.info(f"FEEDBACK | {feedback} | query={query[:80]}")
    except Exception as e:
        _LOGGER.warning(f"Feedback write failed: {e}")


# ===== 渲染函数 =====

def _score_tag(value: float, hi: float = 0.4, lo: float = 0.25) -> str:
    if value > hi:
        return "score-high"
    elif value > lo:
        return "score-mid"
    return "score-low"


def _render_header():
    st.markdown("""
    <div class="app-header">
        <div class="app-header-icon">😎🎶</div>
        <div class="app-header-text">
            <h1>技术支持助手</h1>
            <p>canyoufeelmyworld</p>
        </div>
    </div>
    """, unsafe_allow_html=True)


def _render_empty_state():
    st.markdown("""
    <div class="empty-state">
        <div class="icon">💬</div>
        <div class="title">输入问题，开始查询</div>
    </div>
    """, unsafe_allow_html=True)


def _render_references(evidences: list):
    if not evidences:
        return
    st.markdown('<div class="ref-section">', unsafe_allow_html=True)
    for i, ev in enumerate(evidences):
        meta = ev.get("metadata", {})
        src = meta.get("source_file", "未知来源")
        page = meta.get("page", "")
        section = meta.get("section", "")
        topic = meta.get("topic", "")
        source_level = meta.get("source_level", "")
        text = (ev.get("document", "") or ev.get("text", "")).strip()

        level_map = {"official": "📗", "training": "📘", "inferred": "📙"}
        level_icon = level_map.get(source_level.lower(), "📄")

        st.markdown(f"""
        <div class="ref-item">
            <div class="ref-index">{i + 1}</div>
            <div class="ref-info">
                <div class="ref-src">{level_icon} {src}</div>
                <div class="ref-meta">
                    {f'<span>📄 第{page}页</span>' if page else ''}
                    {f'<span>📂 {section}</span>' if section else ''}
                    {f'<span>🏷️ {topic}</span>' if topic else ''}
                </div>
                <details style="margin-top:6px">
                    <summary style="cursor:pointer; color:#6b7280; font-size:0.82rem">查看证据原文</summary>
                    <div class="ref-evidence-text">{text[:2000] if text else '(无内容)'}</div>
                </details>
            </div>
        </div>
        """, unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)


def _render_retrieval_detail(result: dict):
    evidences = result.get("evidences", [])
    if not evidences:
        return

    timing = result.get("timing", {})
    with st.expander("🔍 检索分析", expanded=False):
        chips = ""
        for label, key in [("分类", "classification"), ("向量检索", "retrieval"),
                            ("Rerank", "rerank"), ("LLM生成", "llm"), ("Grounding", "grounding")]:
            v = timing.get(key, 0)
            if v:
                chips += f'<span class="metric-chip">{label} <span class="val">{v:.2f}s</span></span>'
        st.markdown(f'<div class="metrics-row">{chips}</div>', unsafe_allow_html=True)

        st.markdown(f"*召回 {result.get('retrieval_count', '?')} 条 → Rerank 保留 Top {len(evidences)}*")

        for i, ev in enumerate(evidences):
            meta = ev.get("metadata", {})
            src = meta.get("source_file", "?")[:40]
            section = meta.get("section", "")[:30]
            sim = ev.get("similarity", 0)
            rerank = ev.get("rerank_score", 0)
            text = (ev.get("document", "") or ev.get("text", ""))

            st.markdown(f"""
            <div class="retrieval-item">
                <strong>#{i + 1}</strong> &nbsp;
                <span class="score-tag {_score_tag(rerank, 0.35, 0.24)}">Rerank {rerank:.4f}</span>
                <span class="score-tag {_score_tag(sim, 0.4, 0.25)}">Cos {sim:.4f}</span>
                &nbsp; {src}
                {f'<br><small style="color:#9ca3af">📂 {section}</small>' if section else ''}
            </div>
            """, unsafe_allow_html=True)
            with st.expander(f"证据 #{i + 1} 全文"):
                st.text(text[:1500] if text else "(空)")


def _render_answer_card(msg: dict, msg_id: str):
    result = msg.get("result", {})
    answer = result.get("answer", "")
    is_answerable = result.get("is_answerable", True)
    category = result.get("category", "knowledge")
    grounding = result.get("grounding", {})

    cat_config = {
        "knowledge": ("📖 知识问答", "badge-knowledge"),
        "implementation": ("🔧 实施指导", "badge-implementation"),
        "troubleshooting": ("🔍 故障排查", "badge-troubleshooting"),
    }
    cat_label, cat_css = cat_config.get(category, ("📖 知识问答", "badge-knowledge"))

    refuse_markers = ["没有找到足够依据", "知识库中没有", "未检索到", "没有找到足够依据确认"]
    llm_refused = any(m in answer for m in refuse_markers)

    if llm_refused:
        card_class = "ai-card ai-card-refusal"
        header_badge = '<span class="badge badge-refusal">⚠️ 依据不足</span>'
    else:
        card_class = "ai-card"
        header_badge = f'<span class="badge {cat_css}">{cat_label}</span>'

    timing = result.get("timing", {})
    total_t = timing.get("total", 0)
    time_str = f"{total_t:.1f}s" if total_t else ""
    regenerated = " · 🔄 已重新生成" if result.get("regenerated") else ""

    html = f"""
    <div class="{card_class}">
        <div class="ai-card-header">
            {header_badge}
            <span style="font-size:0.78rem;color:#9ca3af;margin-left:auto">{time_str}{regenerated}</span>
        </div>
        <div class="ai-card-body">
            {_md_to_html(answer)}
        </div>
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)

    if grounding and is_answerable and not llm_refused:
        if not grounding.get("grounded", True):
            unverified = grounding.get("unverified", [])
            if unverified:
                st.markdown(
                    f'<div class="grounding-warn">⚠️ 以上回答中有 <b>{len(unverified)} 处</b>声明在知识库证据中未能充分验证，建议优先确认参考资料原文。</div>',
                    unsafe_allow_html=True,
                )

    st.markdown("#### 📚 参考资料")
    _render_references(result.get("evidences", []))
    _render_retrieval_detail(result)

    c1, c2, c3, c4 = st.columns([0.7, 0.7, 2.6, 1])
    fb_key = f"fb_{msg_id}"

    if st.session_state.feedback_given.get(fb_key):
        c1.markdown("<small>✅ 已反馈</small>", unsafe_allow_html=True)
    else:
        with c1:
            if st.button("👍 有帮助", key=f"hlp_{msg_id}"):
                _log_feedback(result.get("query", ""), answer, "helpful")
                st.session_state.feedback_given[fb_key] = "helpful"
                st.rerun()
        with c2:
            if st.button("👎 没帮助", key=f"nhlp_{msg_id}"):
                _log_feedback(result.get("query", ""), answer, "not_helpful")
                st.session_state.feedback_given[fb_key] = "not_helpful"
                st.rerun()

    with c4:
        if st.button("🔄 重新生成", key=f"regen_{msg_id}"):
            _do_regenerate(msg, msg_id)


def _md_to_html(text: str) -> str:
    """极简 Markdown → HTML（处理标题/列表/引用/加粗 + 证据分级标记高亮）"""
    import re
    lines = text.split("\n")
    out = []
    in_list = False
    in_olist = False

    for line in lines:
        stripped = line.strip()

        if not stripped:
            if in_list:
                out.append("</ul>")
                in_list = False
            if in_olist:
                out.append("</ol>")
                in_olist = False
            out.append("<br>")
            continue

        if stripped.startswith("### "):
            if in_list:
                out.append("</ul>"); in_list = False
            if in_olist:
                out.append("</ol>"); in_olist = False
            out.append(f"<h3>{stripped[4:]}</h3>")
            continue

        if re.match(r"^[-*]\s", stripped):
            if not in_list:
                if in_olist:
                    out.append("</ol>"); in_olist = False
                out.append("<ul>")
                in_list = True
            content = re.sub(r"^[-*]\s+", "", stripped)
            content = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", content)
            out.append(f"<li>{content}</li>")
            continue

        m = re.match(r"^(\d+)\.\s", stripped)
        if m:
            if not in_olist:
                if in_list:
                    out.append("</ul>"); in_list = False
                out.append("<ol>")
                in_olist = True
            content = stripped[m.end():]
            content = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", content)
            out.append(f"<li>{content}</li>")
            continue

        if stripped.startswith("> "):
            if in_list:
                out.append("</ul>"); in_list = False
            if in_olist:
                out.append("</ol>"); in_olist = False
            content = stripped[2:]
            content = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", content)
            out.append(f"<blockquote>{content}</blockquote>")
            continue

        if in_list:
            out.append("</ul>"); in_list = False
        if in_olist:
            out.append("</ol>"); in_olist = False
        line = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", stripped)
        out.append(f"<p>{line}</p>")

    if in_list:
        out.append("</ul>")
    if in_olist:
        out.append("</ol>")

    html = "\n".join(out)

    for tag, css in _EVIDENCE_TAG_STYLES.items():
        html = html.replace(tag, f'<span class="{css}">{tag}</span>')

    return html


def _do_regenerate(msg: dict, msg_id: str):
    query = msg.get("query", "")
    evidences = msg.get("result", {}).get("evidences", [])
    extra_context = msg.get("result", {}).get("extra_context", "")
    chat_history = st.session_state.chat_history[:-2] if len(st.session_state.chat_history) >= 2 else []

    if not evidences:
        st.warning("无法重新生成：未找到原始检索证据。")
        return

    try:
        from src.llm.prompt_builder import build_messages, build_citation_text, ground_citations, classify_query

        with st.spinner("🔄 基于相同证据重新生成中..."):
            classification = classify_query(query)
            messages = build_messages(query, evidences, chat_history if chat_history else None,
                                      extra_context=extra_context)
            pipeline = _load_pipeline()
            answer_text = pipeline.llm_client.chat(messages)
            citations = build_citation_text(evidences)
            grounding = ground_citations(answer_text, evidences)

            msg["result"]["answer"] = answer_text
            msg["result"]["citations"] = citations
            msg["result"]["grounding"] = grounding
            msg["result"]["category"] = classification.category
            msg["result"]["category_confidence"] = classification.confidence
            msg["result"]["regenerated"] = True
            _LOGGER.info(f"REGENERATE | query={query[:80]}")
            st.rerun()
    except Exception as e:
        _LOGGER.error(f"Regenerate failed: {e}", exc_info=True)
        st.error("重新生成失败，请稍后重试。")


# ===== 错误处理 =====
_ERROR_MAP = {
    "DeepSeek": "AI服务暂时不可用，请稍后重试。",
    "chromadb": "知识库检索失败，请稍后重试。",
    "embedding": "模型加载失败，请检查配置。",
    "SentenceTransformer": "模型加载失败，请检查配置。",
    "ConnectionError": "网络连接失败，请检查网络后重试。",
    "Timeout": "请求超时，AI服务响应较慢，请稍后重试。",
}


def _friendly_error(e: Exception) -> str:
    err_str = str(e)
    for key, msg in _ERROR_MAP.items():
        if key.lower() in err_str.lower():
            return msg
    return f"系统异常: {err_str[:200]}"


def _clear_conversation():
    st.session_state.chat_history = []
    st.session_state.chat_messages = []
    st.session_state.last_evidences = None
    st.session_state.last_query = None
    st.session_state.last_result = None
    st.session_state.feedback_given = {}
    _LOGGER.info("Conversation cleared")


# ===== 附件栏（聊天框上方，右侧） =====
def _render_attach_bar():
    """渲染聊天框右下角的附件按钮 + 展开的上传面板"""
    # 附件按钮：靠右对齐
    left, right = st.columns([6, 1])
    with right:
        toggle = st.button("📎", key="attach_toggle", help="拍照 / 文件上传", use_container_width=True)
        if toggle:
            st.session_state.show_attach = not st.session_state.get("show_attach", False)

    if st.session_state.get("show_attach"):
        c1, c2 = st.columns([1, 1])
        with c1:
            photo = st.camera_input("📷 拍照", key="camera_upload")
            if photo is not None:
                st.image(photo, caption="已拍摄照片", use_container_width=True)
                st.caption("⚠️ 纯文本模型暂不支持识别图片，请用文字描述截图关键信息")
        with c2:
            file = st.file_uploader("📁 文件", type=_UPLOAD_TYPES, key="file_upload")
            if file is not None:
                text, kind = _extract_text(file)
                if kind == "text" and text.strip():
                    text = text[:4000]
                    st.session_state["upload_text"] = text
                    st.session_state["upload_meta"] = {"name": file.name, "chars": len(text)}
                    st.success(f"✅ 已加载「{file.name}」（{len(text)} 字）")
                elif kind == "image":
                    st.session_state["upload_text"] = ""
                    st.session_state["upload_meta"] = None
                    st.warning("⚠️ 图片暂不支持识别，请用文字描述")
                else:
                    st.session_state["upload_text"] = ""
                    st.session_state["upload_meta"] = None
                    st.warning("⚠️ 该文件格式暂不支持提取文本")

    # 当前附件状态
    meta = st.session_state.get("upload_meta")
    if meta and st.session_state.get("upload_text"):
        st.caption(f"📎 已附加：{meta['name']}（{meta['chars']} 字）")


# ===== 主入口 =====
def main():
    _render_header()

    # ---- 侧边栏（关于信息） ----
    with st.sidebar:
        st.markdown("### ⚙️ 会话管理")
        if st.button("🗑️ 清空对话历史", use_container_width=True):
            _clear_conversation()
            st.rerun()

        msg_count = len([m for m in st.session_state.chat_messages if m.get("role") == "user"])
        st.caption(f"💬 本轮已提问 {msg_count} 次")

        st.markdown("---")
        st.markdown("### ℹ️ 关于")
        st.caption(
            "**技术支持助手** v0.6\n\n"
            "基于RAG技术检索产品手册、技术培训及实施文档，"
            "为售后工程师提供实施、配置及故障排查参考。\n\n"
            "⚠️ AI回答仅供参考，关键操作请以官方文档为准。"
        )

    # ---- 消息区 ----
    if not st.session_state.chat_messages:
        st.markdown("<br>", unsafe_allow_html=True)
        _render_empty_state()
    else:
        for i, msg in enumerate(st.session_state.chat_messages):
            role = msg.get("role", "")
            msg_id = msg.get("msg_id", str(i))

            if role == "user":
                attachment = msg.get("attachment", "")
                att_html = f'<div class="user-attach">📎 {attachment}</div>' if attachment else ""
                st.markdown(f"""
                <div class="user-msg-wrapper">
                    <div class="user-bubble">
                        <div class="user-label">YOU</div>
                        {msg["content"]}{att_html}
                    </div>
                </div>
                """, unsafe_allow_html=True)
            elif role == "assistant":
                _render_answer_card(msg, msg_id)
                st.markdown("<br>", unsafe_allow_html=True)

    # ---- 附件栏（聊天框上方右侧） ----
    _render_attach_bar()

    # ---- 输入 ----
    query = st.chat_input("输入你的技术问题...")

    if not query or not query.strip():
        return

    query = query.strip()

    if (st.session_state.chat_messages and
            st.session_state.chat_messages[-1].get("content") == query):
        st.warning("请勿重复发送相同问题。")
        return

    extra_context = st.session_state.get("upload_text", "")
    attach_name = (st.session_state.get("upload_meta") or {}).get("name", "")

    msg_id = str(st.session_state.msg_counter)
    st.session_state.msg_counter += 1
    st.session_state.chat_messages.append({
        "role": "user",
        "content": query,
        "msg_id": msg_id + "_u",
        "attachment": attach_name,
    })

    pipeline = _load_pipeline()

    spinner_txt = "🔍 检索知识库并分析中..." if not extra_context else "🔍 检索知识库并分析附件材料..."
    with (st.spinner(spinner_txt)):
        t0 = time.time()
        try:
            result = pipeline.answer(query, extra_context=extra_context)
            wall = time.time() - t0
        except Exception as e:
            _LOGGER.error(f"Pipeline error: {e}", exc_info=True)
            st.error(_friendly_error(e))
            return

    result_dict = {
        "query": result.query,
        "category": result.category,
        "category_confidence": result.category_confidence,
        "is_answerable": result.is_answerable,
        "answer": result.answer,
        "citations": result.citations,
        "evidences": result.evidences,
        "retrieval_count": result.retrieval_count,
        "rerank_count": result.rerank_count,
        "timing": result.timing,
        "grounding": result.grounding,
        "regenerated": False,
        "extra_context": extra_context,
    }

    st.session_state.chat_messages.append({
        "role": "assistant",
        "content": result.answer,
        "query": query,
        "result": result_dict,
        "msg_id": msg_id,
    })
    st.session_state.last_evidences = result.evidences
    st.session_state.last_query = query
    st.session_state.last_result = result_dict
    st.session_state.chat_history.append({"role": "user", "content": query})
    st.session_state.chat_history.append({"role": "assistant", "content": result.answer[:1000]})

    if len(st.session_state.chat_history) > 20:
        st.session_state.chat_history = st.session_state.chat_history[-20:]

    _log_query(query, result_dict, wall)
    st.rerun()


if __name__ == "__main__":
    main()
